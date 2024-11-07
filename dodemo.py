# dodemo.py - A python script to read from excel and write to powerpoint.
#Does not require either excel or powerpoint to be installed ; libraries read / write directly to files

#Python 3.12
#pip 24.0 / 3.12

# Not required but used in the demo
# Visual Studio Code 1.94.2

#Notable python packages -
#pandas - python data analysis library
#pypptx - python powerpoint library, reads and writes .pptx files
#colordict - mostly a dictionary that has predefined colors and english names for them

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from colordict import *
 
# Load the Excel file into a pandas dataframe
excel_file = 'workbook.xlsx'
df = pd.read_excel(excel_file)
print(df)

# Create a PowerPoint presentation object
prs = Presentation()
# This is a little wierd, for any arbitrary ppt presentaiton
# there are layouts that aren't referencable by name
# in a new, blank powerpoint deck, the blank layout is number 6
SLD_LAYOUT_BLANK=6

# Add a slide with a blank
slide_layout = prs.slide_layouts[SLD_LAYOUT_BLANK]
slide = prs.slides.add_slide(slide_layout)

# Define the starting position and dimensions of the shapes
origin_left = Inches(1)
origin_top = Inches(1)
origin_width = Inches(2)
origin_height = Inches(1)

# Define the values that will change during the run
left = origin_left 
top = origin_top 
width = origin_width 
height = origin_height 
hex_dict = ColorDict(mode='hex')
#This dumps out the possible colors by name not required
#but helpful for the demo
print (hex_dict)


#print(MSO_ANCHOR._member_names_)

# Iterate through the dataframe and add rounded rectangles with text to the slide
for i, row in df.iterrows():
    text = row['Desc']
    color = row['Color']
    
    # Add a rounded rectangle shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    top = top + height
    #if the current top is more than 5 inches down the slide:
    #move over by the width + .03 and reset the top
    if (top > Inches(5)): 
       left=left+width+.03    
       top=origin_top

    # Set the fill color of the shape
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_dict[color].lstrip('#'))
    
    # Add text to the shape
    shape.text = text

# Save the presentation
prs.save('output_presentation.pptx')

print("PowerPoint presentation created successfully.")

